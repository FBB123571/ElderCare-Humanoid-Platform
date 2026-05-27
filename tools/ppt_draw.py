"""用 PPT 原生形状绘制框图（非 matplotlib 贴图）。"""
from __future__ import annotations

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

DARK = RGBColor(20, 50, 120)
BLUE = RGBColor(37, 99, 235)
CYAN = RGBColor(14, 165, 233)
GREEN = RGBColor(16, 185, 129)
ORANGE = RGBColor(245, 158, 11)
RED = RGBColor(220, 38, 38)
GRAY = RGBColor(100, 116, 139)
LIGHT = RGBColor(241, 247, 255)
WHITE = RGBColor(255, 255, 255)
PALE_RED = RGBColor(254, 226, 226)
PALE_BLUE = RGBColor(219, 234, 254)


def _box(slide, x, y, w, h, text, fill, *, size=11, bold=True, fg=WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
  s = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
  s.fill.solid()
  s.fill.fore_color.rgb = fill
  s.line.color.rgb = RGBColor(255, 255, 255)
  s.line.width = Pt(1)
  tf = s.text_frame
  tf.clear()
  tf.vertical_anchor = MSO_ANCHOR.MIDDLE
  tf.word_wrap = True
  tf.margin_left = Pt(4)
  tf.margin_right = Pt(4)
  p = tf.paragraphs[0]
  p.text = text
  p.alignment = PP_ALIGN.CENTER
  p.font.name = "微软雅黑"
  p.font.size = Pt(size)
  p.font.bold = bold
  p.font.color.rgb = fg
  return s


def _arrow_h(slide, x1, y, x2, color=GRAY):
  conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y), Inches(x2), Inches(y))
  conn.line.color.rgb = color
  conn.line.width = Pt(2)


def _arrow_v(slide, x, y1, y2, color=GRAY):
  conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x), Inches(y1), Inches(x), Inches(y2))
  conn.line.color.rgb = color
  conn.line.width = Pt(2)


def draw_toc_timeline(slide, left: float, top: float, w: float, h: float):
  """目录时间轴。"""
  items = [
    ("01", "背景与需求", BLUE),
    ("02", "技术与模型", CYAN),
    ("03", "实现与验证", GREEN),
    ("04", "创新与展望", ORANGE),
  ]
  step = h / len(items)
  for i, (num, label, col) in enumerate(items):
    y = top + 0.08 + i * step
    _box(slide, left, y, 0.42, 0.36, num, col, size=12)
    _box(slide, left + 0.5, y, w - 0.55, 0.36, label, LIGHT, size=12, fg=DARK)
    if i < len(items) - 1:
      _arrow_v(slide, left + 0.21, y + 0.38, y + step - 0.06, col)


def draw_architecture(slide, left: float, top: float, w: float, h: float):
  layers = [
    ("交互层", "Web / GUI", CYAN),
    ("认知层", "风险 · 对话 · 规划", BLUE),
    ("感知层", "跌倒 · 情绪 · 视觉", GREEN),
    ("执行层", "机器人 / ROS2", ORANGE),
  ]
  bh = (h - 0.55) / len(layers)
  y = top + 0.05
  cx = left + w / 2
  for title, sub, col in layers:
    _box(slide, left + 0.08, y, w - 0.16, bh - 0.06, f"{title}\n{sub}", col, size=11)
    y += bh
  _box(slide, left + w * 0.22, top + h - 0.48, w * 0.56, 0.42, "CareOrchestrator", DARK, size=10)
  for y0 in [top + bh * i + 0.02 for i in range(1, 4)]:
    _arrow_v(slide, cx, y0 - 0.04, y0 + 0.02)


def draw_state_machine(slide, left: float, top: float, w: float, h: float):
  names = ["监测", "对话", "告警", "紧急"]
  cols = [BLUE, BLUE, BLUE, RED]
  bw = (w - 0.35) / 4
  y = top + h * 0.38
  bh = h * 0.28
  xs = [left + 0.05 + i * (bw + 0.05) for i in range(4)]
  for x, name, col in zip(xs, names, cols):
    _box(slide, x, y, bw, bh, name, col, size=11)
  for i in range(3):
    _arrow_h(slide, xs[i] + bw, y + bh / 2, xs[i + 1], DARK)
  _box(slide, left + 0.1, top + h - 0.42, w - 0.2, 0.32, "紧急态可打断常规任务", PALE_RED, size=10, fg=RED, bold=True)


def draw_pipeline(slide, left: float, top: float, w: float, h: float):
  nodes = ["采集", "特征", "融合", "状态机", "规划", "执行"]
  bw = (w - 0.25) / len(nodes)
  y = top + h * 0.42
  bh = h * 0.22
  x = left + 0.08
  for i, n in enumerate(nodes):
    _box(slide, x, y, bw - 0.04, bh, n, BLUE, size=10)
    if i < len(nodes) - 1:
      _arrow_h(slide, x + bw - 0.04, y + bh / 2, x + bw + 0.02, DARK)
    x += bw


def draw_compare(slide, left: float, top: float, w: float, h: float):
  lw = w * 0.46
  rw = w * 0.46
  gap = w * 0.08
  pairs = [
    ("音箱各自为战", "多模态统一输入"),
    ("摄像头只报警", "风险可解释输出"),
    ("手环难判跌倒", "骨架+规则检测"),
    ("难复现演示", "Web+脚本开源"),
  ]
  row_h = (h - 0.35) / len(pairs)
  y = top + 0.28
  _box(slide, left, top + 0.02, lw, 0.22, "常见做法", GRAY, size=10)
  _box(slide, left + lw + gap, top + 0.02, rw, 0.22, "本作品", DARK, size=10)
  for ltxt, rtxt in pairs:
    _box(slide, left, y, lw, row_h - 0.05, ltxt, PALE_RED, size=10, fg=RED, bold=False)
    _box(slide, left + lw + gap, y, rw, row_h - 0.05, rtxt, PALE_BLUE, size=10, fg=DARK, bold=False)
    y += row_h


def draw_fall_flow(slide, left: float, top: float, w: float, h: float):
  cx = left + w / 2
  steps = [
    (top + 0.05, "骨架输入", CYAN),
    (top + h * 0.28, "算 R、v_y", BLUE),
    (top + h * 0.52, "快速/慢速规则", BLUE),
    (top + h * 0.78, "触发紧急链路", RED),
  ]
  bw, bh = w * 0.62, h * 0.16
  x = cx - bw / 2
  for y, txt, col in steps:
    _box(slide, x, y, bw, bh, txt, col, size=11)
  for i in range(len(steps) - 1):
    _arrow_v(slide, cx, steps[i][0] + bh, steps[i + 1][0] - 0.02)


def draw_emergency_chain(slide, left: float, top: float, w: float, h: float):
  steps = [("告警", ORANGE), ("语音", BLUE), ("通知", RED), ("手势", GREEN)]
  bw = (w - 0.2) / 4
  y = top + h * 0.4
  bh = h * 0.2
  x = left + 0.06
  for label, col in steps:
    _box(slide, x, y, bw - 0.05, bh, label, col, size=11)
    _arrow_h(slide, x + bw - 0.05, y + bh / 2, x + bw, DARK)
    x += bw


def draw_risk_formula(slide, left: float, top: float, w: float, h: float):
  _box(slide, left + 0.05, top + 0.05, w - 0.1, 0.38, "S = 0.45·跌倒 + 0.25·情绪 + 0.30·健康", LIGHT, size=13, fg=DARK)
  items = [("跌倒", GREEN), ("情绪", CYAN), ("健康", ORANGE)]
  bw = (w - 0.2) / 3
  y = top + h * 0.48
  x = left + 0.08
  for label, col in items:
    _box(slide, x, y, bw - 0.06, h * 0.38, label, col, size=12)
    x += bw
  _box(slide, left + 0.1, top + h - 0.32, w - 0.2, 0.26, "输出分数 + 等级 + 原因", DARK, size=10)


def draw_aging_bars(slide, left: float, top: float, w: float, h: float):
  """竖向柱状图（形状绘制）。"""
  years = ["2020", "2025", "2030", "2035"]
  vals = [13.5, 16.5, 20.5, 24.0]
  cols = [BLUE, CYAN, ORANGE, RED]
  chart_h = h * 0.72
  base_y = top + h * 0.88
  bar_w = w * 0.16
  gap = (w - bar_w * 4) / 5
  x = left + gap
  for year, v, col in zip(years, vals, cols):
    bh = chart_h * (v / 28)
    y = base_y - bh
    _box(slide, x, y, bar_w, bh, f"{v}%", col, size=10)
    tb = slide.shapes.add_textbox(Inches(x), Inches(base_y + 0.04), Inches(bar_w), Inches(0.2))
    p = tb.text_frame.paragraphs[0]
    p.text = year
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "微软雅黑"
    p.font.size = Pt(9)
    p.font.color.rgb = GRAY
    x += bar_w + gap
  cap = slide.shapes.add_textbox(Inches(left), Inches(top + 0.02), Inches(w), Inches(0.22))
  cap.text_frame.paragraphs[0].text = "60岁及以上人口占比（%）"
  cap.text_frame.paragraphs[0].font.name = "微软雅黑"
  cap.text_frame.paragraphs[0].font.size = Pt(10)
  cap.text_frame.paragraphs[0].font.color.rgb = DARK


def draw_metrics(slide, left: float, top: float, w: float, h: float):
  labels = ["Precision", "Recall", "F1", "紧急触发"]
  row_h = h / 5
  y = top + 0.1
  for lab in labels:
    tb = slide.shapes.add_textbox(Inches(left), Inches(y), Inches(1.1), Inches(row_h * 0.7))
    p = tb.text_frame.paragraphs[0]
    p.text = lab
    p.font.name = "微软雅黑"
    p.font.size = Pt(10)
    p.font.color.rgb = DARK
    _box(slide, left + 1.15, y + 0.02, w - 1.25, row_h * 0.65, "100%", GREEN, size=10)
    y += row_h


def draw_roadmap(slide, left: float, top: float, w: float, h: float):
  phases = [("已完成", "Web·评测·仓库", GREEN), ("进行中", "软著·真机", BLUE), ("计划", "数据集·场景", ORANGE)]
  pw = (w - 0.15) / 3
  y = top + h * 0.25
  ph = h * 0.5
  x = left + 0.04
  for title, body, col in phases:
    _box(slide, x, y, pw - 0.04, ph, f"{title}\n{body}", col, size=11)
    if x + pw < left + w - 0.1:
      _arrow_h(slide, x + pw - 0.04, y + ph / 2, x + pw + 0.02, GRAY)
    x += pw


def draw_innovation_grid(slide, left: float, top: float, w: float, h: float):
  tags = [
    ("融合", "多源信号加权"), ("闭环", "紧急可打断"),
    ("视觉", "MediaPipe"), ("工程", "评测脚本"),
    ("场景", "居家养老"),
  ]
  cols = 2
  cw = (w - 0.12) / cols
  rh = (h - 0.1) / 3
  for i, (t, d) in enumerate(tags):
    r, c = divmod(i, cols)
    x = left + 0.04 + c * (cw + 0.04)
    y = top + 0.05 + r * rh
    _box(slide, x, y, cw, rh - 0.06, f"{t}\n{d}", BLUE if i % 2 == 0 else CYAN, size=10)


def draw_llm_flow(slide, left: float, top: float, w: float, h: float):
  bw = w * 0.28
  y = top + h * 0.35
  bh = h * 0.22
  x1, x2, x3 = left + 0.04, left + 0.36, left + 0.68
  _box(slide, x1, y, bw, bh, "用户\n情绪", CYAN, size=10)
  _box(slide, x2, y, bw, bh, "对话\n模块", DARK, size=10)
  _box(slide, x3, y, bw, bh, "安抚\n话术", GREEN, size=10)
  _arrow_h(slide, x1 + bw, y + bh / 2, x2, DARK)
  _arrow_h(slide, x2 + bw, y + bh / 2, x3, DARK)
  _box(slide, left + 0.12, top + h * 0.72, w - 0.24, 0.2, "默认 Mock，可换 API", LIGHT, size=9, fg=DARK, bold=False)
