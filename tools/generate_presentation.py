#!/usr/bin/env python3
"""生成国赛风格答辩 PPT：docs/答辩_CareCompanion.pptx"""
from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

import sys
sys.path.insert(0, str(Path(__file__).resolve().parents[1]))
from tools.ppt_layout import (
  CAPTION_H,
  COL_L,
  COL_L_W,
  COL_R,
  COL_R_W,
  CONTENT_BOTTOM,
  CONTENT_H,
  CONTENT_TOP,
  FIG_BODY_H,
  FIG_CAP_H,
  FIG_LEFT,
  FIG_TOP,
  FIG_TOTAL_H,
  FIG_WIDTH,
  FOOTER_H,
  GUTTER,
  HEADER_H,
  IMG_PAD,
  MARGIN,
  RIBBON_H,
  SLIDE_H,
  SLIDE_W,
  TITLE_H,
  TITLE_TOP,
  body_height,
  ribbon_top,
)

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs" / "答辩_CareCompanion.pptx"
ASSETS = ROOT / "docs" / "assets"
SCI = ASSETS / "sci"
DIAG = ASSETS / "diagrams"

WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(28, 32, 44)
DARK = RGBColor(20, 50, 120)
GRAY = RGBColor(90, 98, 115)
LIGHT_BLUE = RGBColor(241, 247, 255)
ACCENT = RGBColor(245, 158, 11)
SOFT_BLUE = RGBColor(219, 234, 254)
BORDER = RGBColor(186, 200, 220)
FONT_CN = "微软雅黑"

COMP_LINE1 = "第二十八届中国机器人及人工智能大赛"
COMP_LINE2 = "创新赛 · 机器人创新赛道"

TEAM_NAME = "从容应队"
SCHOOL = "中山大学"
PROJECT_TITLE = "基于大语言模型的智能养老陪伴机器人仿真平台"
TEAM_ID = "CRAIC2026-TEAM-8FJQKLMI"
PROJECT_ID = "CRAIC20264ZEFT1DF"
CAPTAIN = "刘小凡"
MEMBER = "白冉"
ADVISOR = "彭键清"
GITHUB_URL = "https://github.com/FBB123571/ElderCare-Humanoid-Platform"


def _slide(prs: Presentation):
  return prs.slides.add_slide(prs.slide_layouts[6])


def _set_run_font(run, size: int, *, bold=False, color=BLACK, name=FONT_CN):
  run.font.name = name
  run.font.size = Pt(size)
  run.font.bold = bold
  if color:
    run.font.color.rgb = color
  try:
    rpr = run._r.rPr
    if rpr is not None:
      from lxml import etree
      rfonts = rpr.find(qn("w:rFonts"))
      if rfonts is None:
        rfonts = etree.SubElement(rpr, qn("w:rFonts"))
      rfonts.set(qn("w:eastAsia"), name)
      rfonts.set(qn("w:ascii"), name)
      rfonts.set(qn("w:hAnsi"), name)
  except Exception:
    pass


def _set_para_font(p, size: int, *, bold=False, color=BLACK):
  if not p.runs:
    p.text = p.text or " "
  for run in p.runs:
    _set_run_font(run, size, bold=bold, color=color)
  if not p.runs:
    _set_run_font(p.add_run(), size, bold=bold, color=color)


def _bg_white(slide):
  slide.background.fill.solid()
  slide.background.fill.fore_color.rgb = WHITE


def _header_bar(slide):
  bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(SLIDE_W), Inches(HEADER_H))
  bar.fill.solid()
  bar.fill.fore_color.rgb = DARK
  bar.line.fill.background()
  box = slide.shapes.add_textbox(Inches(MARGIN), Inches(0.06), Inches(SLIDE_W - 2 * MARGIN), Inches(0.38))
  p = box.text_frame.paragraphs[0]
  p.text = f"{COMP_LINE1}  |  {COMP_LINE2}"
  _set_para_font(p, 9, bold=True, color=WHITE)


def _footer_bar(slide, hint: str = ""):
  y = SLIDE_H - FOOTER_H
  bar = slide.shapes.add_shape(1, Inches(0), Inches(y), Inches(SLIDE_W), Inches(FOOTER_H))
  bar.fill.solid()
  bar.fill.fore_color.rgb = SOFT_BLUE
  bar.line.fill.background()
  txt = f"{TEAM_NAME} · {SCHOOL} · {CAPTAIN}/{MEMBER} · 指导教师 {ADVISOR}"
  if hint:
    txt = f"{hint}  |  {txt}"
  box = slide.shapes.add_textbox(Inches(MARGIN), Inches(y + 0.05), Inches(SLIDE_W - 2 * MARGIN), Inches(0.2))
  p = box.text_frame.paragraphs[0]
  p.text = txt
  p.alignment = PP_ALIGN.CENTER
  _set_para_font(p, 8, color=GRAY)


def _title(slide, text: str, *, center=False):
  left = MARGIN if not center else MARGIN
  width = SLIDE_W - 2 * MARGIN
  box = slide.shapes.add_textbox(Inches(left), Inches(TITLE_TOP), Inches(width), Inches(TITLE_H))
  p = box.text_frame.paragraphs[0]
  p.text = text
  _set_para_font(p, 20, bold=True, color=DARK)
  p.alignment = PP_ALIGN.CENTER if center else PP_ALIGN.LEFT
  if not center:
    line = slide.shapes.add_shape(1, Inches(left), Inches(TITLE_TOP + TITLE_H - 0.04), Inches(1.25), Inches(0.055))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()


def _ribbon(slide, text: str):
  y = ribbon_top()
  bar = slide.shapes.add_shape(1, Inches(COL_L), Inches(y), Inches(COL_L_W), Inches(RIBBON_H))
  bar.fill.solid()
  bar.fill.fore_color.rgb = DARK
  bar.line.fill.background()
  box = slide.shapes.add_textbox(Inches(COL_L + 0.12), Inches(y + 0.07), Inches(COL_L_W - 0.24), Inches(RIBBON_H - 0.1))
  p = box.text_frame.paragraphs[0]
  p.text = text
  p.alignment = PP_ALIGN.CENTER
  _set_para_font(p, 11, bold=True, color=WHITE)


def _style_cell(cell, text: str, *, bg=DARK, fg=WHITE, bold=True, size=11):
  cell.text = ""
  p = cell.text_frame.paragraphs[0]
  p.text = text
  p.alignment = PP_ALIGN.CENTER
  _set_para_font(p, size, bold=bold, color=fg)
  cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
  cell.fill.solid()
  cell.fill.fore_color.rgb = bg
  cell.margin_left = Pt(4)
  cell.margin_right = Pt(4)
  cell.margin_top = Pt(2)
  cell.margin_bottom = Pt(2)


def _bullets_panel(
  slide,
  items: list[str],
  *,
  highlights: list[int] | None = None,
  size: int = 12,
  takeaway: str = "",
):
  """左栏要点：每条独立文本框 + 顶对齐，避免垂直居中造成行距错乱。"""
  highlights = highlights or []
  top = CONTENT_TOP
  height = CONTENT_H
  foot_h = 0.34 if takeaway else 0.0
  pad_top, pad_side = 0.12, 0.18

  panel = slide.shapes.add_shape(1, Inches(COL_L), Inches(top), Inches(COL_L_W), Inches(height))
  panel.fill.solid()
  panel.fill.fore_color.rgb = LIGHT_BLUE
  panel.line.color.rgb = BORDER
  panel.line.width = Pt(0.8)

  stripe = slide.shapes.add_shape(1, Inches(COL_L), Inches(top), Inches(0.08), Inches(height))
  stripe.fill.solid()
  stripe.fill.fore_color.rgb = DARK
  stripe.line.fill.background()

  text_top = top + pad_top
  text_h = height - pad_top - 0.12 - foot_h
  n = max(1, len(items))
  slot_h = text_h / n

  for i, line in enumerate(items):
    y = text_top + i * slot_h
    box = slide.shapes.add_textbox(
      Inches(COL_L + pad_side), Inches(y),
      Inches(COL_L_W - pad_side * 2), Inches(slot_h - 0.02),
    )
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Pt(2)
    tf.margin_top = Pt(0)
    p = tf.paragraphs[0]
    p.text = line if line.startswith("●") else f"● {line}"
    _set_para_font(p, size, bold=(i in highlights), color=DARK if i in highlights else BLACK)
    p.line_spacing = 1.12
    p.space_after = Pt(2)

  if takeaway:
    fy = top + height - foot_h
    foot = slide.shapes.add_shape(
      1, Inches(COL_L + 0.10), Inches(fy), Inches(COL_L_W - 0.20), Inches(foot_h - 0.02),
    )
    foot.fill.solid()
    foot.fill.fore_color.rgb = DARK
    foot.line.fill.background()
    fb = slide.shapes.add_textbox(
      Inches(COL_L + 0.14), Inches(fy + 0.07), Inches(COL_L_W - 0.28), Inches(foot_h - 0.14),
    )
    tf = fb.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    fp = tf.paragraphs[0]
    fp.text = takeaway
    fp.alignment = PP_ALIGN.CENTER
    _set_para_font(fp, 11, bold=True, color=WHITE)


def _fit_image(path: Path, max_w: float, max_h: float) -> tuple[float, float]:
  with Image.open(path) as im:
    px_w, px_h = im.size
  aspect = px_h / px_w
  w, h = max_w, max_w * aspect
  if h > max_h:
    h, w = max_h, max_h / aspect
  return w, h


def _draw_figure_panel(slide, *, left: float, top: float, width: float, height: float) -> None:
  """右栏配图底板，与左栏要点区对齐、不越界。"""
  panel = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
  panel.fill.solid()
  panel.fill.fore_color.rgb = WHITE
  panel.line.color.rgb = BORDER
  panel.line.width = Pt(0.75)


def _place_figure(
  slide,
  path: Path,
  *,
  left: float,
  top: float,
  max_w: float,
  max_h: float,
  caption: str = "",
  fill_ratio: float = 0.96,
  framed: bool = True,
) -> None:
  """在矩形区域内等比缩放（contain），禁止溢出到左栏或页眉页脚。"""
  if not path.exists():
    return

  cap_h = FIG_CAP_H if caption else 0.0
  region_h = max_h
  img_h_max = max(0.5, region_h - cap_h - 0.06)
  box_w = max(0.5, max_w * fill_ratio)
  box_h = max(0.5, img_h_max * fill_ratio)

  w, h = _fit_image(path, box_w, box_h)
  # 二次钳制，防止浮点误差越界
  if w > max_w:
    h *= max_w / w
    w = max_w
  if h > img_h_max:
    w *= img_h_max / h
    h = img_h_max

  x = left + (max_w - w) / 2
  y = top + (img_h_max - h) / 2

  if framed:
    pad = 0.04
    frame = slide.shapes.add_shape(
      1, Inches(x - pad), Inches(y - pad),
      Inches(min(max_w, w + 2 * pad)), Inches(min(img_h_max, h + 2 * pad)),
    )
    frame.fill.solid()
    frame.fill.fore_color.rgb = RGBColor(252, 253, 255)
    frame.line.color.rgb = BORDER
    frame.line.width = Pt(0.6)

  slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))

  if caption:
    cap_top = top + region_h - cap_h
    cap = slide.shapes.add_textbox(Inches(left), Inches(cap_top), Inches(max_w), Inches(cap_h))
    tf = cap.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    cp = tf.paragraphs[0]
    cp.text = caption
    cp.alignment = PP_ALIGN.CENTER
    _set_para_font(cp, 8, color=GRAY)


def _place_right_figure(slide, path: Path, *, caption: str = "") -> None:
  """标准左文右图页的右栏配图（图体与图注分区）。"""
  panel_left = COL_R + 0.04
  panel_w = COL_R_W - 0.08
  _draw_figure_panel(slide, left=panel_left, top=FIG_TOP - 0.02, width=panel_w, height=FIG_TOTAL_H)
  sep = slide.shapes.add_shape(1, Inches(COL_L + COL_L_W + GUTTER * 0.45), Inches(CONTENT_TOP),
                               Inches(0.012), Inches(CONTENT_H))
  sep.fill.solid()
  sep.fill.fore_color.rgb = BORDER
  sep.line.fill.background()
  _place_figure(
    slide, path,
    left=FIG_LEFT, top=FIG_TOP, max_w=FIG_WIDTH, max_h=FIG_BODY_H + FIG_CAP_H,
    caption=caption, fill_ratio=0.92, framed=True,
  )


def _content_slide(prs, title: str, footer_hint: str = ""):
  slide = _slide(prs)
  _bg_white(slide)
  _header_bar(slide)
  _footer_bar(slide, footer_hint)
  _title(slide, title)
  return slide


def _caption_for(kind: str, custom: str) -> str:
  if custom:
    return custom
  return {
    "diagram": "结构示意图",
    "chart": "实验结果",
    "photo": "系统实拍",
    "sci": "论文风格配图",
  }.get(kind, "")


def _info_cards_column(slide, left: float, width: float, rows: list[tuple[str, str]]):
  """右栏参赛信息表：单元格垂直居中、行高均分。"""
  top = CONTENT_TOP
  height = CONTENT_H
  n = len(rows)
  tbl = slide.shapes.add_table(n, 2, Inches(left), Inches(top), Inches(width), Inches(height)).table
  label_w = Inches(0.92)
  tbl.columns[0].width = label_w
  tbl.columns[1].width = Inches(width) - label_w

  for r, (label, value) in enumerate(rows):
    _style_cell(tbl.cell(r, 0), label, bg=DARK, fg=WHITE, bold=True, size=10)
    c1 = tbl.cell(r, 1)
    c1.text = ""
    p = c1.text_frame.paragraphs[0]
    p.text = value
    p.alignment = PP_ALIGN.LEFT
    fs = 10 if len(value) <= 28 else 8
    _set_para_font(p, fs, bold=False, color=DARK)
    c1.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    c1.text_frame.word_wrap = True
    c1.fill.solid()
    c1.fill.fore_color.rgb = LIGHT_BLUE
    c1.margin_left = Pt(6)
    c1.margin_right = Pt(4)


def _bullets_plain(slide, items: list[str], *, size: int = 12):
  """左栏编号条目（白底卡片，与答辩提纲蓝底面板区分）。"""
  y = CONTENT_TOP + 0.10
  for i, line in enumerate(items):
    chip = slide.shapes.add_shape(1, Inches(COL_L), Inches(y), Inches(COL_L_W), Inches(0.54))
    chip.fill.solid()
    chip.fill.fore_color.rgb = WHITE
    chip.line.color.rgb = BORDER
    num = slide.shapes.add_shape(9, Inches(COL_L + 0.12), Inches(y + 0.12), Inches(0.28), Inches(0.28))
    num.fill.solid()
    num.fill.fore_color.rgb = ACCENT if i == 0 else DARK
    num.line.fill.background()
    nb = slide.shapes.add_textbox(Inches(COL_L + 0.12), Inches(y + 0.15), Inches(0.28), Inches(0.22))
    np = nb.text_frame.paragraphs[0]
    np.text = str(i + 1)
    np.alignment = PP_ALIGN.CENTER
    _set_para_font(np, 9, bold=True, color=WHITE)
    tb = slide.shapes.add_textbox(Inches(COL_L + 0.48), Inches(y + 0.14), Inches(COL_L_W - 0.58), Inches(0.34))
    p = tb.text_frame.paragraphs[0]
    p.text = line
    p.word_wrap = True
    _set_para_font(p, size, color=DARK)
    y += 0.60


def _slide_lr_dual_fig(
  prs,
  title: str,
  bullets: list[str],
  fig_top: str,
  fig_bottom: str,
  *,
  cap_top: str = "",
  cap_bottom: str = "",
  takeaway: str = "",
):
  slide = _content_slide(prs, title)
  _bullets_panel(slide, bullets, takeaway=takeaway, size=12)
  _draw_figure_panel(slide, left=COL_R + 0.04, top=FIG_TOP - 0.02, width=COL_R_W - 0.08, height=FIG_TOTAL_H)
  gap = 0.10
  half = (FIG_BODY_H - gap) / 2
  top_y = FIG_TOP
  for fig_name, cap, y, h in [
    (fig_top, cap_top, top_y, half),
    (fig_bottom, cap_bottom, top_y + half + gap, half),
  ]:
    path = SCI / fig_name
    if path.exists():
      _place_figure(
        slide, path, left=FIG_LEFT, top=y, max_w=FIG_WIDTH, max_h=h,
        caption=cap, fill_ratio=0.93, framed=False,
      )
  return slide


def _slide_team_info(prs):
  slide = _content_slide(prs, "参赛信息")
  _bullets_plain(slide, [
    "赛项：创新赛 · 机器人创新赛道",
    "作品：智能养老陪伴机器人仿真平台",
    "特色：理论清晰 · 系统可演示 · 指标可复现",
  ], size=11)

  panel = slide.shapes.add_shape(1, Inches(COL_R), Inches(CONTENT_TOP), Inches(COL_R_W), Inches(CONTENT_H))
  panel.fill.solid()
  panel.fill.fore_color.rgb = WHITE
  panel.line.color.rgb = BORDER

  dash = ASSETS / "ppt_full_dashboard.png"
  if dash.exists():
    _place_figure(
      slide, dash, left=COL_R + 0.10, top=CONTENT_TOP + 0.08,
      max_w=COL_R_W - 0.20, max_h=1.55, caption="Web 控制台预览", fill_ratio=0.96, framed=False,
    )

  tbl_top = CONTENT_TOP + 1.72
  tbl_h = CONTENT_H - 1.80
  rows = [
    ("团队", TEAM_NAME),
    ("学校", SCHOOL),
    ("队长 / 队员", f"{CAPTAIN} / {MEMBER}"),
    ("指导教师", ADVISOR),
    ("团队编号", TEAM_ID),
    ("作品编号", PROJECT_ID),
  ]
  tbl = slide.shapes.add_table(len(rows), 2, Inches(COL_R + 0.08), Inches(tbl_top),
                             Inches(COL_R_W - 0.16), Inches(tbl_h)).table
  tbl.columns[0].width = Inches(0.88)
  tbl.columns[1].width = Inches(COL_R_W - 0.16 - 0.88)
  for r, (label, value) in enumerate(rows):
    _style_cell(tbl.cell(r, 0), label, bg=DARK, fg=WHITE, bold=True, size=9)
    c1 = tbl.cell(r, 1)
    c1.text = ""
    p = c1.text_frame.paragraphs[0]
    p.text = value
    p.alignment = PP_ALIGN.LEFT
    fs = 8 if len(value) > 24 else 9
    _set_para_font(p, fs, color=DARK)
    c1.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    c1.text_frame.word_wrap = True
    c1.fill.solid()
    c1.fill.fore_color.rgb = LIGHT_BLUE if r % 2 == 0 else WHITE
  return slide


def _slide_lr(
  prs,
  title: str,
  bullets: list[str],
  image: str | Path,
  *,
  diagram: bool = True,
  sci: bool = False,
  figure_kind: str = "",
  bullet_size: int = 13,
  caption: str = "",
  highlight_last: bool = False,
  footer_hint: str = "",
  takeaway: str = "",
):
  slide = _content_slide(prs, title, footer_hint)
  hi = [len(bullets) - 1] if highlight_last and bullets else []
  _bullets_panel(slide, bullets, highlights=hi, size=bullet_size, takeaway=takeaway)

  if sci:
    sub = SCI / image
  elif diagram:
    sub = DIAG / image
    if not sub.exists():
      sub = ASSETS / image
  else:
    sub = ASSETS / image

  kind = figure_kind or ("sci" if sci else ("diagram" if diagram else "photo"))
  cap = _caption_for(kind, caption)
  _place_right_figure(slide, sub, caption=cap)
  return slide


def _slide_draw(prs, title: str, bullets: list[str], draw_fn, *, takeaway: str = "", bullet_size: int = 12, footer_hint: str = ""):
  slide = _content_slide(prs, title, footer_hint)
  _bullets_panel(slide, bullets, takeaway=takeaway, size=bullet_size)
  draw_fn(slide, COL_R + 0.05, CONTENT_TOP + 0.04, COL_R_W - 0.1, CONTENT_H - 0.08)
  return slide


def _slide_skeleton_annotated(prs):
  """骨架页：左图 + 右侧文字注释（不重复贴风险面板图）。"""
  slide = _content_slide(prs, "骨架分析（上传照片）", "系统实拍")
  _bullets_panel(slide, [
    "MediaPipe 提取 33 个骨架关键点",
    "由肩-髋-踝计算宽高比 R",
    "竖直速度突增触发跌倒规则",
    "上传照片与摄像头共用同一逻辑",
  ], size=12)
  _draw_figure_panel(slide, left=COL_R + 0.04, top=FIG_TOP - 0.02, width=COL_R_W - 0.08, height=FIG_TOTAL_H)
  sep = slide.shapes.add_shape(1, Inches(COL_L + COL_L_W + GUTTER * 0.45), Inches(CONTENT_TOP),
                               Inches(0.012), Inches(CONTENT_H))
  sep.fill.solid()
  sep.fill.fore_color.rgb = BORDER
  sep.line.fill.background()

  pose = ASSETS / "ppt_mediapipe_pose.png"
  img_left = FIG_LEFT + 0.06
  img_w = 2.75
  img_top = FIG_TOP + 0.08
  img_h = FIG_BODY_H - 0.10
  if pose.exists():
    _place_figure(
      slide, pose, left=img_left, top=img_top, max_w=img_w, max_h=img_h,
      caption="MediaPipe 骨架叠加", fill_ratio=0.97, framed=False,
    )

  ann_left = img_left + img_w + 0.16
  ann_w = FIG_LEFT + FIG_WIDTH - ann_left - 0.06
  notes = [
    ("① 关键点", "鼻、肩、髋、踝等构成人体构型"),
    ("② 宽高比 R", "躯干投影宽高比异常 → 姿态变化"),
    ("③ 竖直速度 vy", "关键点整体下移超阈 → 判定跌倒"),
  ]
  y = img_top + 0.12
  for title, desc in notes:
    card = slide.shapes.add_shape(1, Inches(ann_left), Inches(y), Inches(ann_w), Inches(0.88))
    card.fill.solid()
    card.fill.fore_color.rgb = LIGHT_BLUE
    card.line.color.rgb = BORDER
    card.line.width = Pt(0.6)
    tb = slide.shapes.add_textbox(Inches(ann_left + 0.14), Inches(y + 0.10), Inches(ann_w - 0.24), Inches(0.70))
    tf = tb.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    p0.text = title
    _set_para_font(p0, 12, bold=True, color=ACCENT)
    p1 = tf.add_paragraph()
    p1.text = desc
    _set_para_font(p1, 11, color=DARK)
    p1.space_before = Pt(3)
    y += 0.98
  return slide


def _screenshot_duo(prs, title: str, left_img: str, right_img: str, cap_l: str = "", cap_r: str = ""):
  slide = _content_slide(prs, title, "系统实拍")
  gap = 0.14
  half = (SLIDE_W - 2 * MARGIN - gap) / 2
  region_h = CONTENT_H - 0.08
  top = CONTENT_TOP + 0.04
  for i, (img, cap) in enumerate([(left_img, cap_l), (right_img, cap_r)]):
    left = MARGIN + i * (half + gap)
    _draw_figure_panel(slide, left=left, top=top, width=half, height=region_h)
    _place_figure(
      slide, ASSETS / img,
      left=left + 0.08, top=top + 0.06,
      max_w=half - 0.16, max_h=region_h - 0.08,
      caption=cap, fill_ratio=0.94, framed=False,
    )
  return slide


def _slide_open_source(prs):
  """开源页：左表格 + 右二维码与文字清单（无叠图）。"""
  slide = _content_slide(prs, "开源与可复现性")

  rows = [
    ("GitHub", GITHUB_URL),
    ("团队", f"{TEAM_NAME}（{SCHOOL}）"),
    ("本地演示", "bash scripts/run_web.sh"),
    ("演示录像", "docs/assets/demo_carecompanion.mp4"),
    ("软著", "申请中"),
  ]
  n = len(rows)
  tbl = slide.shapes.add_table(n, 2, Inches(COL_L), Inches(CONTENT_TOP), Inches(COL_L_W), Inches(CONTENT_H)).table
  tbl.columns[0].width = Inches(0.95)
  tbl.columns[1].width = Inches(COL_L_W - 0.95)
  for r, (label, value) in enumerate(rows):
    _style_cell(tbl.cell(r, 0), label, bg=DARK, fg=WHITE, bold=True, size=9)
    c1 = tbl.cell(r, 1)
    c1.text = ""
    p = c1.text_frame.paragraphs[0]
    p.text = value
    p.alignment = PP_ALIGN.LEFT
    fs = 8 if len(value) > 36 else 9
    _set_para_font(p, fs, color=DARK)
    c1.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    c1.text_frame.word_wrap = True
    c1.fill.solid()
    c1.fill.fore_color.rgb = LIGHT_BLUE if r % 2 == 0 else WHITE
    c1.margin_left = Pt(5)

  rx = COL_R + 0.06
  rw = COL_R_W - 0.12
  panel = slide.shapes.add_shape(1, Inches(rx), Inches(CONTENT_TOP), Inches(rw), Inches(CONTENT_H))
  panel.fill.solid()
  panel.fill.fore_color.rgb = WHITE
  panel.line.color.rgb = BORDER

  qr = ASSETS / "ppt_demo_qr.png"
  qr_size = 1.55
  if qr.exists():
    _place_figure(
      slide, qr, left=rx + (rw - qr_size) / 2, top=CONTENT_TOP + 0.12,
      max_w=qr_size, max_h=qr_size, caption="", fill_ratio=1.0, framed=False,
    )
  cap = slide.shapes.add_textbox(Inches(rx), Inches(CONTENT_TOP + 0.12 + qr_size + 0.04), Inches(rw), Inches(0.22))
  cp = cap.text_frame.paragraphs[0]
  cp.text = "扫码访问 GitHub 仓库"
  cp.alignment = PP_ALIGN.CENTER
  _set_para_font(cp, 9, bold=True, color=DARK)

  checks = [
    "✓  完整源码与 README",
    "✓  一键启动脚本 run_web.sh",
    "✓  fall_eval_report.json",
    "✓  pytest 单元测试",
    "○  软著申请中",
  ]
  box = slide.shapes.add_textbox(Inches(rx + 0.15), Inches(CONTENT_TOP + 2.05), Inches(rw - 0.3), Inches(1.55))
  tf = box.text_frame
  tf.word_wrap = True
  for i, line in enumerate(checks):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = line
    col = RGBColor(22, 101, 52) if line.startswith("✓") else RGBColor(180, 83, 9)
    _set_para_font(p, 10, color=col)
    p.space_after = Pt(5)


def _hero_cover(prs):
  slide = _slide(prs)
  _header_bar(slide)
  sh = SLIDE_H

  left_w = 5.0
  left_bg = slide.shapes.add_shape(1, Inches(0), Inches(HEADER_H), Inches(left_w), Inches(sh - HEADER_H))
  left_bg.fill.solid()
  left_bg.fill.fore_color.rgb = LIGHT_BLUE
  left_bg.line.fill.background()

  hero = ASSETS / "ppt_mediapipe_pose.png"
  if not hero.exists():
    hero = ASSETS / "ppt_full_dashboard.png"
  if hero.exists():
    area_h = sh - HEADER_H - 0.25
    _place_figure(slide, hero, left=0.15, top=HEADER_H + 0.12, max_w=left_w - 0.3, max_h=area_h,
                  caption="MediaPipe 骨架 · 可现场演示", fill_ratio=0.98)

  rx = left_w
  right_bg = slide.shapes.add_shape(1, Inches(rx), Inches(HEADER_H), Inches(SLIDE_W - rx), Inches(sh - HEADER_H))
  right_bg.fill.solid()
  right_bg.fill.fore_color.rgb = DARK
  right_bg.line.fill.background()
  deco = slide.shapes.add_shape(1, Inches(rx), Inches(HEADER_H), Inches(0.1), Inches(sh - HEADER_H))
  deco.fill.solid()
  deco.fill.fore_color.rgb = ACCENT
  deco.line.fill.background()

  texts = [
    ("CareCompanion", 28, WHITE, True),
    ("智能养老陪伴机器人仿真平台", 14, RGBColor(147, 197, 253), True),
    ("感知 · 决策 · 执行  一体化闭环", 12, RGBColor(226, 232, 240), False),
  ]
  ty = HEADER_H + 0.42
  for text, sz, col, bold in texts:
    box = slide.shapes.add_textbox(Inches(rx + 0.22), Inches(ty), Inches(4.55), Inches(0.46))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    _set_para_font(p, sz, bold=bold, color=col)
    ty += 0.50 if sz > 20 else 0.36

  card_y = ty + 0.28
  card = slide.shapes.add_shape(1, Inches(rx + 0.22), Inches(card_y), Inches(4.4), Inches(1.52))
  card.fill.solid()
  card.fill.fore_color.rgb = RGBColor(30, 64, 130)
  card.line.color.rgb = RGBColor(100, 140, 200)

  infos = [
    f"团队：{TEAM_NAME}（{SCHOOL}）",
    f"队员：{CAPTAIN}（队长）· {MEMBER}",
    f"指导教师：{ADVISOR}",
    f"作品编号：{PROJECT_ID}",
  ]
  box = slide.shapes.add_textbox(Inches(rx + 0.36), Inches(card_y + 0.12), Inches(4.1), Inches(1.28))
  tf = box.text_frame
  tf.word_wrap = True
  for i, line in enumerate(infos):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = line
    _set_para_font(p, 10, color=RGBColor(220, 230, 255))
    p.space_after = Pt(4)

  tag = slide.shapes.add_textbox(Inches(rx + 0.22), Inches(card_y + 1.58), Inches(4.4), Inches(0.32))
  tp = tag.text_frame.paragraphs[0]
  tp.text = "基于大语言模型的居家养老陪伴仿真平台"
  tp.alignment = PP_ALIGN.CENTER
  _set_para_font(tp, 8, color=RGBColor(148, 163, 184))


def _part_divider(prs, part: str, subtitle: str, preview: list[str], sci_thumb: str = "", asset_thumb: str = ""):
  slide = _slide(prs)
  _bg_white(slide)
  _header_bar(slide)
  _footer_bar(slide)

  left_w = 4.75
  block_h = CONTENT_BOTTOM - HEADER_H
  main = slide.shapes.add_shape(1, Inches(0), Inches(HEADER_H), Inches(left_w), Inches(block_h))
  main.fill.solid()
  main.fill.fore_color.rgb = DARK
  main.line.fill.background()

  num = part.replace("第", "").replace("部分", "")[:1]
  circle = slide.shapes.add_shape(9, Inches(0.42), Inches(HEADER_H + 0.35), Inches(0.72), Inches(0.72))
  circle.fill.solid()
  circle.fill.fore_color.rgb = ACCENT
  circle.line.fill.background()
  nb = slide.shapes.add_textbox(Inches(0.42), Inches(HEADER_H + 0.42), Inches(0.72), Inches(0.45))
  np = nb.text_frame.paragraphs[0]
  np.text = num
  np.alignment = PP_ALIGN.CENTER
  _set_para_font(np, 22, bold=True, color=WHITE)

  ty = HEADER_H + 1.15
  box = slide.shapes.add_textbox(Inches(0.38), Inches(ty), Inches(4.2), Inches(0.5))
  p = box.text_frame.paragraphs[0]
  p.text = part
  _set_para_font(p, 24, bold=True, color=WHITE)

  sub = slide.shapes.add_textbox(Inches(0.4), Inches(ty + 0.52), Inches(4.2), Inches(0.38))
  sp = sub.text_frame.paragraphs[0]
  sp.text = subtitle
  _set_para_font(sp, 12, color=RGBColor(191, 219, 254))

  chip_y = ty + 1.05
  for line in preview:
    chip = slide.shapes.add_shape(1, Inches(0.42), Inches(chip_y), Inches(4.0), Inches(0.42))
    chip.fill.solid()
    chip.fill.fore_color.rgb = RGBColor(30, 64, 130)
    chip.line.color.rgb = RGBColor(100, 140, 200)
    cb = slide.shapes.add_textbox(Inches(0.55), Inches(chip_y + 0.08), Inches(3.75), Inches(0.28))
    cp = cb.text_frame.paragraphs[0]
    cp.text = line
    _set_para_font(cp, 11, color=RGBColor(226, 232, 240))
    chip_y += 0.48

  right_w = SLIDE_W - left_w
  right_bg = slide.shapes.add_shape(1, Inches(left_w), Inches(HEADER_H), Inches(right_w), Inches(block_h))
  right_bg.fill.solid()
  right_bg.fill.fore_color.rgb = LIGHT_BLUE
  right_bg.line.fill.background()
  if asset_thumb:
    thumb_path = ASSETS / asset_thumb
  elif sci_thumb:
    thumb_path = SCI / sci_thumb
  else:
    thumb_path = None
  if thumb_path and thumb_path.exists():
    pad = 0.18
    _place_figure(
      slide, thumb_path,
      left=left_w + pad, top=CONTENT_TOP + 0.12,
      max_w=right_w - 2 * pad, max_h=CONTENT_H - 0.22,
      caption="", fill_ratio=0.90, framed=True,
    )


def _screenshot_slide(prs, title: str, image_file: str, caption: str = "", *, large: bool = False):
  slide = _content_slide(prs, title, "系统实拍")
  path = ASSETS / image_file
  left, w = MARGIN, SLIDE_W - 2 * MARGIN
  top, h = CONTENT_TOP + 0.02, CONTENT_H - 0.04
  pad = 0.04 if large else 0.10
  _draw_figure_panel(slide, left=left, top=top, width=w, height=h)
  _place_figure(
    slide, path,
    left=left + pad, top=top + pad * 0.5,
    max_w=w - 2 * pad, max_h=h - 0.04,
    caption=caption, fill_ratio=0.98 if large else 0.94, framed=False,
  )


def _thanks_slide(prs):
  """致谢页：居中排版，底部仪表盘 + 二维码并列。"""
  slide = _slide(prs)
  _header_bar(slide)
  _footer_bar(slide)
  sh = SLIDE_H

  body = slide.shapes.add_shape(1, Inches(0), Inches(HEADER_H), Inches(SLIDE_W), Inches(sh - HEADER_H - FOOTER_H))
  body.fill.solid()
  body.fill.fore_color.rgb = LIGHT_BLUE
  body.line.fill.background()

  cx = SLIDE_W / 2
  title_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.05), Inches(SLIDE_W - 1.6), Inches(0.75))
  tp = title_box.text_frame.paragraphs[0]
  tp.text = "敬请批评指正"
  tp.alignment = PP_ALIGN.CENTER
  _set_para_font(tp, 32, bold=True, color=DARK)

  line = slide.shapes.add_shape(1, Inches(cx - 0.9), Inches(1.82), Inches(1.8), Inches(0.045))
  line.fill.solid()
  line.fill.fore_color.rgb = ACCENT
  line.line.fill.background()

  sub = slide.shapes.add_textbox(Inches(1.2), Inches(2.0), Inches(SLIDE_W - 2.4), Inches(0.95))
  tf = sub.text_frame
  for i, line_txt in enumerate([
    f"{TEAM_NAME} · {SCHOOL}",
    f"队长 {CAPTAIN}  ·  队员 {MEMBER}",
    f"指导教师 {ADVISOR}",
  ]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = line_txt
    p.alignment = PP_ALIGN.CENTER
    _set_para_font(p, 13 if i == 0 else 12, bold=(i == 0), color=DARK if i == 0 else GRAY)
    p.space_after = Pt(4)

  url_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.95), Inches(SLIDE_W - 3.0), Inches(0.28))
  up = url_box.text_frame.paragraphs[0]
  up.text = GITHUB_URL
  up.alignment = PP_ALIGN.CENTER
  _set_para_font(up, 9, color=DARK)

  row_y = 3.45
  thumb_w = 3.35
  qr_s = 1.35
  gap = 0.35
  total_w = thumb_w + gap + qr_s
  x0 = cx - total_w / 2

  thumb = ASSETS / "ppt_full_dashboard.png"
  if thumb.exists():
    _place_figure(
      slide, thumb, left=x0, top=row_y, max_w=thumb_w, max_h=1.55,
      caption="Web 控制台", fill_ratio=0.96, framed=True,
    )
  qr = ASSETS / "ppt_demo_qr.png"
  if qr.exists():
    _place_figure(
      slide, qr, left=x0 + thumb_w + gap, top=row_y + 0.1,
      max_w=qr_s, max_h=qr_s, caption="GitHub", fill_ratio=1.0, framed=False,
    )


def _slide_qr_dual(prs):
  slide = _content_slide(prs, "扫码查看代码与演示视频")
  _bullets_panel(slide, [
    "GitHub 完整源码与文档",
    "在线 Web 演示（现场可复现）",
    "演示录像已附仓库",
    "欢迎评委扫码查阅",
  ], size=13)
  _draw_figure_panel(slide, left=COL_R + 0.04, top=FIG_TOP - 0.02, width=COL_R_W - 0.08, height=FIG_TOTAL_H)
  qr = ASSETS / "ppt_demo_qr.png"
  demo = ASSETS / "ppt_header_perception.png"
  if not demo.exists():
    demo = ASSETS / "ppt_full_dashboard.png"
  qr_w = 1.72
  gap = 0.14
  img_top = FIG_TOP + 0.06
  fig_body = FIG_BODY_H - 0.18
  if qr.exists():
    _place_figure(
      slide, qr, left=FIG_LEFT, top=img_top + 0.22,
      max_w=qr_w, max_h=qr_w + 0.15, caption="仓库二维码", fill_ratio=0.90, framed=False,
    )
  if demo.exists():
    _place_figure(
      slide, demo, left=FIG_LEFT + qr_w + gap, top=img_top,
      max_w=FIG_WIDTH - qr_w - gap, max_h=fig_body,
      caption="Web 感知界面", fill_ratio=0.88, framed=False,
    )


def _ensure_qr_code() -> None:
  qr = ASSETS / "ppt_demo_qr.png"
  if qr.exists():
    return
  try:
    import qrcode
    img = qrcode.make(GITHUB_URL, box_size=8, border=2)
    ASSETS.mkdir(parents=True, exist_ok=True)
    img.save(qr)
  except ImportError:
    pass


def main():
  import sys
  sys.path.insert(0, str(ROOT))
  from tools import ppt_sci_figures as sci

  _ensure_qr_code()
  sci.gen_all()
  try:
    from tools.clean_pose_asset import clean_pose_image, SRC as POSE_SRC, OUT as POSE_OUT

    tmp_pose = POSE_SRC.with_suffix(".tmp.png")
    if clean_pose_image(POSE_SRC, tmp_pose):
      tmp_pose.replace(POSE_OUT)
  except Exception:
    pass

  prs = Presentation()
  prs.slide_width = Inches(SLIDE_W)
  prs.slide_height = Inches(SLIDE_H)

  _hero_cover(prs)

  _slide_lr(prs, "答辩提纲", [
    "一、为什么做：老龄化与跌倒风险",
    "二、怎么做：模型、架构与算法",
    "三、做出来什么：系统演示与测试",
    "四、有什么特色：创新点与后续",
  ], "outline_timeline.png", sci=True, figure_kind="chart", caption="答辩结构时间轴")

  _slide_team_info(prs)

  _part_divider(prs, "第一部分", "背景 · 需求 · 定位",
                ["老龄化", "方案对比", "项目目标"], asset_thumb="ppt_emotion_input.png")

  _slide_lr(prs, "社会背景：老龄化与跌倒风险", [
    "独居老人变多，跌倒了没人及时发现很常见",
    "监测设备多，但很少把告警和机器人动作连起来",
    "政策在推智慧养老，我们按居家场景做验证",
    "用仿真平台先把流程跑通，再对接人形机",
  ], "aging_trend.png", sci=True, figure_kind="chart", caption="Fig. 老龄化趋势（公开数据示意）",
     takeaway="难点：发现慢，闭环弱")

  _slide_lr(prs, "和常见方案比一比", [
    "我们不是把音箱、摄像头、手环简单拼在一起",
    "用统一调度把感知、决策、执行串成一条链",
    "风险结果能解释，方便答辩和后期调参",
    "仓库里带 Web 和脚本，老师可以现场点着看",
  ], "compare_table.png", sci=True, figure_kind="chart", caption="方案对比表")

  _slide_lr(prs, "项目要做什么", [
    "CareCompanion：养老场景的人形陪伴软件",
    "核心是 CareOrchestrator，负责模块调度",
    "跌倒要能检、要能触发紧急流程",
    "对话和机器人动作要能联动",
  ], "project_goals_combo.png", sci=True, bullet_size=13,
     caption="架构与状态机总览", takeaway="目标在系统里都能点到")

  _part_divider(prs, "第二部分", "模型 · 架构 · 算法",
                ["四层架构", "状态机", "融合与检测"], sci_thumb="risk_fusion.png")

  _slide_lr(prs, "系统怎么分层", [
    "最上面是 Web/GUI，老师和评委直接看",
    "中间做风险、对话和任务规划",
    "下面接跌倒、情绪、视觉等感知",
    "最底层通过适配器连仿真或 ROS2",
  ], "architecture_layers.png", sci=True, caption="四层架构详图")

  _slide_lr(prs, "状态机怎么切换", [
    "平时监测，对话时走倾诉流程",
    "风险升高先告警，确认跌倒进紧急",
    "紧急状态可以打断正在进行的对话",
    "每次循环都会记下分数、回复和动作",
  ], "state_machine.png", sci=True, caption="有限状态机")

  _slide_lr(prs, "风险分数怎么算", [
    "跌倒、情绪、健康三路各占一部分权重",
    "不是端到端黑盒，会给出原因列表",
    "阈值写在配置里，不同场景能改",
    "方便和医生、护理人员沟通",
  ], "risk_fusion.png", sci=True, figure_kind="chart", caption="可解释加权融合")

  _slide_lr(prs, "跌倒怎么判断", [
    "用 MediaPipe 提骨架，算宽高比和速度",
    "快速跌倒和慢速躺倒两套规则",
    "在 6 组合成场景上 F1 到 100%",
    "支持摄像头，也支持上传照片",
  ], "fall_decision_tree.png", sci=True, caption="规则判定树")

  _slide_lr(prs, "对话模块", [
    "先识别老人话里的情绪",
    "再由 DialogueManager 生成回复",
    "默认用 Mock，现场不依赖外网",
    "需要时可以换成 LLM 接口",
  ], "dialogue_sequence.png", sci=True, caption="对话时序图")

  _slide_lr(prs, "数据怎么往下走", [
    "一帧感知数据进来先提特征",
    "融合后交给状态机判断阶段",
    "规划器决定机器人做什么动作",
    "Web 页面上能看到整条链路",
  ], "pipeline_swimlane.png", sci=True, caption="单帧泳道数据流")

  _slide_lr(prs, "确认跌倒之后怎么办", [
    "先播告警，再语音安抚",
    "然后通知家属，最后做举手等动作",
    "EmergencyNotifier 里留了记录",
    "脚本测过端到端触发率 100%",
  ], "emergency_timing.png", sci=True, figure_kind="chart", caption="紧急响应时序",
     takeaway="紧急流程在演示里能完整走通")

  _part_divider(prs, "第三部分", "实现 · 测试 · 演示",
                ["Web", "视觉", "指标"], asset_thumb="ppt_full_dashboard.png")

  _slide_lr(prs, "Web 端长什么样", [
    "后端 FastAPI，前端浏览器打开即可",
    "一页里能看到风险、对话和指令",
    "可以上传照片做骨架分析",
    "内置剧本：日常—倾诉—跌倒",
  ], "ppt_full_dashboard.png", diagram=False, caption="主界面")

  _slide_skeleton_annotated(prs)
  _screenshot_duo(prs, "风险与情绪输入", "ppt_risk_panel.png", "ppt_emotion_input.png", "风险面板", "情绪输入")
  _screenshot_duo(prs, "紧急过程记录", "ppt_chat_emergency.png", "ppt_robot_commands.png", "对话日志", "机器人指令")
  _screenshot_slide(prs, "紧急时感知页", "ppt_header_perception.png", "状态切到紧急后的感知区")

  _slide_lr(prs, "测试结果", [
    "跌倒检测 P/R/F1：合成 6 场景均为 100%",
    "紧急链路脚本跑通率 100%",
    "单帧决策在 CPU 上 < 50ms（见右下图）",
    "单元测试和无头演示都过了",
  ], "eval_results.png", sci=True, figure_kind="chart", caption="P/R/F1 均为 100%（见左栏）",
     takeaway="数据：fall_eval_report.json")

  _slide_lr(prs, "性能与延迟", [
    "骨架提取约占时最多，仍满足实时要求",
    "融合与规则判定开销较小",
    "端到端在 CPU 上约 42ms",
    "现场演示不依赖 GPU",
  ], "latency_bar.png", sci=True, figure_kind="chart", caption="单帧链路耗时（本机粗测）")

  _slide_open_source(prs)

  _part_divider(prs, "第四部分", "创新 · 计划", ["五个方向", "真机", "软著"], asset_thumb="ppt_robot_commands.png")

  _slide_lr(prs, "创新点", [
    "多路信号融合，结果能解释",
    "紧急态在状态机里优先级最高",
    "视觉用 MediaPipe，不是纯参数滑块",
    "代码和评测脚本都开源",
    "场景对准居家养老",
  ], "innovation_radar.png", sci=True, figure_kind="chart", caption="创新维度自评")

  _slide_lr(prs, "后面打算做什么", [
    "对接 Unitree 等人形真机",
    "补充公开跌倒数据集",
    "做多房间、弱光等复杂场景",
    "去养老院做访谈和试用",
    "软著和专利在准备",
  ], "roadmap_gantt.png", sci=True, figure_kind="chart", caption="里程碑甘特图（示意）")

  slide = _content_slide(prs, "参考资料")
  box = slide.shapes.add_textbox(Inches(0.55), Inches(1.35), Inches(8.9), Inches(3.5))
  tf = box.text_frame
  for i, line in enumerate([
    "[1] Google MediaPipe Pose 技术文档",
    "[2] Unitree / ROS2 开发文档",
    "[3] 「十四五」国家老龄事业发展规划",
    "[4] 项目技术报告 docs/TECHNICAL_REPORT.md",
    f"[5] 源码 {GITHUB_URL}",
  ]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = line
    _set_para_font(p, 13)
    p.space_after = Pt(10)

  _slide_qr_dual(prs)
  _thanks_slide(prs)

  OUT.parent.mkdir(parents=True, exist_ok=True)
  prs.save(str(OUT))
  print(f"✅ 已生成: {OUT}  共 {len(prs.slides)} 页")


if __name__ == "__main__":
  main()
